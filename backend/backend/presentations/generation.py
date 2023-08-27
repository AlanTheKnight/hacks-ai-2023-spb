from __future__ import annotations

from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from random import randint
import os


PRESENTATION_WIDTH = Inches(16)
PRESENTATION_HEIGHT = Inches(9)


def add_custom_logo(config, slide: Slide):
    if config["logo"].get("path") is not None:
        if os.path.isfile(config["logo"]["path"]):
            border_offset = Inches(0.25) + Inches(config["logo"]["size"])
            slide.shapes.add_picture(
                config["logo"]["path"],
                PRESENTATION_WIDTH - border_offset,
                PRESENTATION_HEIGHT - border_offset,
                width=Inches(config["logo"]["size"]),
            )


def apply_text_formatting(
    config, obj, color: str = "000000", bold=False, is_title=False
):
    obj.text_frame.paragraphs[0].font.name = config["font"]["name"]
    obj.text_frame.paragraphs[0].font.bold = bold
    obj.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(color)
    font_size = (
        config["font"]["title_size"] if is_title else config["font"]["regular_size"]
    )
    obj.text_frame.paragraphs[0].font.size = Pt(font_size)


def set_slide_background(config, slide: Slide):
    background = slide.background
    fill = background.fill
    if config["bg"]["type"] == "gradient":
        fill.gradient()
        fill.gradient_stops[0].color.rgb = RGBColor.from_string(config["bg"]["color1"])
        fill.gradient_stops[1].color.rgb = RGBColor.from_string(config["bg"]["color2"])
        fill.gradient_angle = randint(0, 360)
    elif config["bg"]["type"] == "solid":
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(config["bg"]["color"])


def create_title_slide(config, slide: Slide):
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = config["name"]
    title.width = PRESENTATION_WIDTH - Inches(2)
    title.left = Inches(1)
    title.top = Inches(3.5)

    subtitle.text = config["pptx_data"]["about"]
    subtitle.width = Inches(14)
    subtitle.left = Inches(1)
    subtitle.top = Inches(3.5) + Inches(1)

    apply_text_formatting(config, title, bold=True, is_title=True)
    apply_text_formatting(config, subtitle)


def generate_title_text_slide(config, prs, title_value, text_value):
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = title_value
    title.width = PRESENTATION_WIDTH - Inches(2)
    title.left = Inches(1)
    title.top = Inches(0.8)
    apply_text_formatting(config, title, bold=True, is_title=True)

    if text_value is not None:
        text = slide.placeholders[1]
        text.text = text_value
        text.width = PRESENTATION_WIDTH - Inches(2)
        text.left = Inches(1)
        text.top = Inches(3)
        apply_text_formatting(config, text)

    set_slide_background(config, slide)
    add_custom_logo(config, slide)


def driver(config, output_path: str):
    prs = Presentation()

    prs.slide_width = PRESENTATION_WIDTH
    prs.slide_height = PRESENTATION_HEIGHT

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    create_title_slide(config, slide)
    set_slide_background(config, slide)
    add_custom_logo(config, slide)

    generate_title_text_slide(config, prs, "Проблема", config["pptx_data"]["problem"])
    generate_title_text_slide(
        config, prs, "Наше решение", config["pptx_data"]["solution"]
    )
    generate_title_text_slide(
        config, prs, "Целевая аудитория", config["pptx_data"]["target"]
    )
    generate_title_text_slide(config, prs, "Наша цель", config["pptx_data"]["goal"])
    generate_title_text_slide(config, prs, "Задачи", config["pptx_data"]["activity"])
    generate_title_text_slide(
        config, prs, "Преймущества", config["pptx_data"]["advantages"]
    )
    generate_title_text_slide(
        config, prs, "Удобство", config["pptx_data"]["convenience"]
    )

    generate_title_text_slide(
        config,
        prs,
        "Трекшн и финансы",
        "Трекшн, партнерства, выручка, количество клиентов, CAC - LTV",
    )
    generate_title_text_slide(
        config,
        prs,
        "Бизнес-модель",
        "Бизнес-модель стартапа, тарифы, условия для клиентов",
    )
    generate_title_text_slide(config, prs, "Наша команда", None)
    generate_title_text_slide(config, prs, "Инвестиции", None)
    generate_title_text_slide(config, prs, "Roadmap", None)
    generate_title_text_slide(config, prs, "Контакты", None)

    prs.save(output_path)
